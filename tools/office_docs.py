"""
office_docs.py — Word & PowerPoint Tools
==========================================
Covers: .docx (read, create, update) and .pptx (read, create slides)
Approach:
  1. Download file from OneDrive via Microsoft Graph
  2. Process locally with python-docx / python-pptx
  3. Upload the modified file back to OneDrive
Auth: reuses ms365.get_access_token() — same MSAL flow
"""

import os
import io
import tempfile
import requests
from pathlib import Path
from typing import Optional

# ── python-docx ─────────────────────────────────────────────────────────────
try:
    from docx import Document
    from docx.shared import Pt, RGBColor, Inches
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

# ── python-pptx ─────────────────────────────────────────────────────────────
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt as PptPt, Emu
    from pptx.dml.color import RGBColor as PptRGB
    from pptx.enum.text import PP_ALIGN
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


GRAPH_BASE = "https://graph.microsoft.com/v1.0"


# ── Shared OneDrive helpers ───────────────────────────────────────────────────

def _get_token() -> str:
    from tools.ms365 import get_access_token
    return get_access_token()


def _find_onedrive_item(filename: str) -> dict:
    """Search OneDrive for a file by name. Returns the full item dict."""
    token = _get_token()
    headers = {"Authorization": f"Bearer {token}"}
    resp = requests.get(
        f"{GRAPH_BASE}/me/drive/root/search(q='{filename}')"
        f"?$select=id,name,webUrl,size",
        headers=headers,
    )
    resp.raise_for_status()
    items = resp.json().get("value", [])
    for item in items:
        if item["name"].lower() == filename.lower():
            return item
    if items:
        return items[0]
    raise FileNotFoundError(f"File not found in OneDrive: {filename}")


def _download_to_bytes(item_id: str) -> bytes:
    """Download a OneDrive file by item ID, return raw bytes."""
    token = _get_token()
    headers = {"Authorization": f"Bearer {token}"}
    resp = requests.get(
        f"{GRAPH_BASE}/me/drive/items/{item_id}/content",
        headers=headers, allow_redirects=True,
    )
    resp.raise_for_status()
    return resp.content


def _upload_bytes(data: bytes, filename: str, folder: str = "/") -> dict:
    """Upload bytes to OneDrive, overwriting if the file exists."""
    token = _get_token()
    headers = {
        "Authorization": f"Bearer {token}",
        "Content-Type": "application/octet-stream",
    }
    folder = folder.rstrip("/")
    dest = f"root:{folder}/{filename}:" if folder else f"root:/{filename}:"
    resp = requests.put(
        f"{GRAPH_BASE}/me/drive/{dest}/content",
        headers=headers, data=data,
    )
    resp.raise_for_status()
    result = resp.json()
    return {"id": result.get("id"), "webUrl": result.get("webUrl"), "name": result.get("name")}


# ══════════════════════════════════════════════════════════════════════════════
# WORD (.docx)
# ══════════════════════════════════════════════════════════════════════════════

def _require_docx():
    if not DOCX_AVAILABLE:
        raise ImportError(
            "python-docx is not installed.\n"
            "Run:  pip install python-docx"
        )


def read_word_document(filename: str) -> dict:
    """
    Read a Word document from OneDrive and return its text content.

    Args:
        filename: Name of the .docx file e.g. 'Report.docx'

    Returns:
        {
          "filename": str,
          "sections": [{"heading": str, "content": str}, ...],
          "full_text": str,
          "paragraph_count": int,
          "webUrl": str
        }
    """
    _require_docx()
    item = _find_onedrive_item(filename)
    raw = _download_to_bytes(item["id"])

    doc = Document(io.BytesIO(raw))
    sections = []
    current_heading = "Document"
    current_content = []

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        if para.style.name.startswith("Heading"):
            if current_content:
                sections.append({"heading": current_heading, "content": "\n".join(current_content)})
                current_content = []
            current_heading = text
        else:
            current_content.append(text)

    if current_content:
        sections.append({"heading": current_heading, "content": "\n".join(current_content)})

    full_text = "\n\n".join(
        f"## {s['heading']}\n{s['content']}" for s in sections
    )

    return {
        "filename": item["name"],
        "sections": sections,
        "full_text": full_text,
        "paragraph_count": len(doc.paragraphs),
        "webUrl": item.get("webUrl", ""),
    }


def create_word_document(
    filename: str,
    title: str,
    sections: list[dict],
    upload_folder: str = "/",
) -> dict:
    """
    Create a new Word document and upload it to OneDrive.

    Args:
        filename:      Output filename e.g. 'My Report.docx'
        title:         Document title (shown as H1 at the top)
        sections:      List of {"heading": str, "content": str} dicts
                       Content supports plain text. Use \\n for paragraphs.
        upload_folder: OneDrive folder to upload to (default: root)

    Returns:
        {"status": "created", "filename": ..., "webUrl": ..., "id": ...}
    """
    _require_docx()
    doc = Document()

    # Document title
    title_para = doc.add_heading(title, level=0)
    title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER

    doc.add_paragraph()  # Spacer

    for section in sections:
        heading = section.get("heading", "")
        content = section.get("content", "")

        if heading:
            doc.add_heading(heading, level=1)

        for line in content.split("\n"):
            line = line.strip()
            if not line:
                continue
            if line.startswith("- ") or line.startswith("• "):
                doc.add_paragraph(line[2:], style="List Bullet")
            elif line[0].isdigit() and ". " in line[:4]:
                doc.add_paragraph(line.split(". ", 1)[1], style="List Number")
            else:
                doc.add_paragraph(line)

    # Save to bytes
    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)

    # Upload
    result = _upload_bytes(buf.read(), filename, upload_folder)
    return {"status": "created", "filename": filename, **result}


def update_word_document(
    filename: str,
    append_sections: list[dict] = None,
    replace_paragraph: dict = None,
) -> dict:
    """
    Update an existing Word document on OneDrive.

    Args:
        filename:          .docx filename in OneDrive
        append_sections:   List of {"heading": str, "content": str} to add at the end
        replace_paragraph: {"find": "old text", "replace": "new text"} — replaces first match

    Returns:
        {"status": "updated", "filename": ..., "webUrl": ...}
    """
    _require_docx()
    item = _find_onedrive_item(filename)
    raw = _download_to_bytes(item["id"])
    doc = Document(io.BytesIO(raw))

    # Replace paragraph text
    if replace_paragraph:
        find_text = replace_paragraph.get("find", "")
        new_text = replace_paragraph.get("replace", "")
        for para in doc.paragraphs:
            if find_text.lower() in para.text.lower():
                # Preserve run formatting, just replace text
                for run in para.runs:
                    if find_text.lower() in run.text.lower():
                        run.text = run.text.replace(find_text, new_text)
                break

    # Append sections
    if append_sections:
        doc.add_page_break()
        for section in append_sections:
            heading = section.get("heading", "")
            content = section.get("content", "")
            if heading:
                doc.add_heading(heading, level=1)
            for line in content.split("\n"):
                line = line.strip()
                if line:
                    doc.add_paragraph(line)

    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)

    result = _upload_bytes(buf.read(), filename)
    return {"status": "updated", "filename": filename, **result}


def list_word_headings(filename: str) -> list[str]:
    """
    Return just the heading structure of a Word document (like a table of contents).

    Returns:
        List of heading strings with level indicators e.g. ["# Introduction", "## Background"]
    """
    _require_docx()
    item = _find_onedrive_item(filename)
    raw = _download_to_bytes(item["id"])
    doc = Document(io.BytesIO(raw))

    headings = []
    for para in doc.paragraphs:
        if para.style.name.startswith("Heading"):
            level = int(para.style.name.split()[-1]) if para.style.name[-1].isdigit() else 1
            prefix = "#" * level
            headings.append(f"{prefix} {para.text.strip()}")
    return headings


# ══════════════════════════════════════════════════════════════════════════════
# POWERPOINT (.pptx)
# ══════════════════════════════════════════════════════════════════════════════

def _require_pptx():
    if not PPTX_AVAILABLE:
        raise ImportError(
            "python-pptx is not installed.\n"
            "Run:  pip install python-pptx"
        )


def read_presentation(filename: str) -> dict:
    """
    Read a PowerPoint file from OneDrive and extract text from all slides.

    Args:
        filename: .pptx filename e.g. 'Quarterly Review.pptx'

    Returns:
        {
          "filename": str,
          "slide_count": int,
          "slides": [{"slide_number": int, "title": str, "content": [str], "notes": str}],
          "webUrl": str
        }
    """
    _require_pptx()
    item = _find_onedrive_item(filename)
    raw = _download_to_bytes(item["id"])

    prs = Presentation(io.BytesIO(raw))
    slides = []

    for i, slide in enumerate(prs.slides, 1):
        title = ""
        content_lines = []
        notes = ""

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue
            if shape.shape_type == 13:  # Picture
                continue
            if hasattr(shape, "placeholder_format") and shape.placeholder_format:
                ph_idx = shape.placeholder_format.idx
                if ph_idx == 0:   # Title
                    title = text
                    continue
            content_lines.append(text)

        # Notes
        if slide.has_notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            if notes_frame:
                notes = notes_frame.text.strip()

        slides.append({
            "slide_number": i,
            "title": title,
            "content": content_lines,
            "notes": notes,
        })

    return {
        "filename": item["name"],
        "slide_count": len(prs.slides),
        "slides": slides,
        "webUrl": item.get("webUrl", ""),
    }


def create_presentation(
    filename: str,
    title: str,
    slides: list[dict],
    theme: str = "dark",
    upload_folder: str = "/",
) -> dict:
    """
    Create a new PowerPoint presentation and upload it to OneDrive.

    Args:
        filename:      Output filename e.g. 'Q3 Review.pptx'
        title:         Presentation title (shown on cover slide)
        slides:        List of slide dicts:
                         {"title": str, "bullets": [str], "notes": str (optional)}
        theme:         'dark' (default) or 'light'
        upload_folder: OneDrive folder path

    Returns:
        {"status": "created", "filename": ..., "slide_count": int, "webUrl": ...}
    """
    _require_pptx()
    from pptx.util import Inches, Pt as PPt
    from pptx.dml.color import RGBColor as PRGB

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Colour scheme
    if theme == "dark":
        bg_colour  = PRGB(0x1a, 0x1a, 0x2e)
        title_colour = PRGB(0x64, 0xff, 0xda)
        text_colour  = PRGB(0xe0, 0xe0, 0xe0)
        accent_colour = PRGB(0xe9, 0x45, 0x60)
    else:
        bg_colour  = PRGB(0xff, 0xff, 0xff)
        title_colour = PRGB(0x0f, 0x34, 0x60)
        text_colour  = PRGB(0x1a, 0x1a, 0x2e)
        accent_colour = PRGB(0xe9, 0x45, 0x60)

    def _set_bg(slide, colour):
        """Set slide background colour."""
        from pptx.util import Pt
        from pptx.oxml.ns import qn
        from lxml import etree
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = colour

    def _add_textbox(slide, text, left, top, width, height, font_size, colour, bold=False, align="left"):
        from pptx.util import Inches, Pt, Emu
        txBox = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        run = p.runs[0]
        run.font.size = Pt(font_size)
        run.font.color.rgb = colour
        run.font.bold = bold
        if align == "center":
            p.alignment = PP_ALIGN.CENTER
        elif align == "right":
            p.alignment = PP_ALIGN.RIGHT

    # ── Cover slide ──────────────────────────────────────────────────────────
    cover_layout = prs.slide_layouts[6]  # Blank
    cover = prs.slides.add_slide(cover_layout)
    _set_bg(cover, bg_colour)

    # Accent bar at top
    from pptx.util import Inches, Pt, Emu
    bar = cover.shapes.add_shape(
        1,  # Rectangle
        Inches(0), Inches(0), Inches(13.33), Inches(0.08)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent_colour
    bar.line.fill.background()

    _add_textbox(cover, title,
                 left=1, top=2.5, width=11.33, height=1.5,
                 font_size=40, colour=title_colour, bold=True, align="center")
    _add_textbox(cover, f"{len(slides)} slides",
                 left=1, top=4.3, width=11.33, height=0.5,
                 font_size=16, colour=PRGB(0x88, 0x92, 0xb0), align="center")

    # ── Content slides ────────────────────────────────────────────────────────
    for slide_data in slides:
        slide_title   = slide_data.get("title", "")
        bullets       = slide_data.get("bullets", [])
        notes_text    = slide_data.get("notes", "")

        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        _set_bg(slide, bg_colour)

        # Accent bar
        bar2 = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(0.06))
        bar2.fill.solid()
        bar2.fill.fore_color.rgb = accent_colour
        bar2.line.fill.background()

        # Slide title
        _add_textbox(slide, slide_title,
                     left=0.5, top=0.3, width=12.33, height=0.9,
                     font_size=28, colour=title_colour, bold=True)

        # Divider line
        line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.35), Inches(12.33), Inches(0.02))
        line.fill.solid()
        line.fill.fore_color.rgb = PRGB(0x2a, 0x3a, 0x5e)
        line.line.fill.background()

        # Bullets
        if bullets:
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.33), Inches(5.5))
            tf = txBox.text_frame
            tf.word_wrap = True
            for j, bullet in enumerate(bullets):
                p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                p.text = f"▸  {bullet}"
                run = p.runs[0]
                run.font.size = Pt(18)
                run.font.color.rgb = text_colour
                p.space_before = Pt(8)

        # Speaker notes
        if notes_text:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes_text

    # Save and upload
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)

    result = _upload_bytes(buf.read(), filename, upload_folder)
    return {
        "status": "created",
        "filename": filename,
        "slide_count": len(slides) + 1,  # +1 for cover
        **result,
    }


def add_slide_to_presentation(
    filename: str,
    slide_title: str,
    bullets: list[str],
    notes: str = "",
    position: int = None,
) -> dict:
    """
    Add a new slide to an existing PowerPoint on OneDrive.

    Args:
        filename:    .pptx filename in OneDrive
        slide_title: Title for the new slide
        bullets:     List of bullet point strings
        notes:       Speaker notes (optional)
        position:    Slide position (default: append at end)

    Returns:
        {"status": "updated", "slide_count": n, "webUrl": ...}
    """
    _require_pptx()
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor as PRGB

    item = _find_onedrive_item(filename)
    raw = _download_to_bytes(item["id"])
    prs = Presentation(io.BytesIO(raw))

    # Use a blank layout
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.9))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = slide_title
    p.runs[0].font.size = Pt(28)
    p.runs[0].font.bold = True

    # Bullets
    if bullets:
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5.5))
        tf2 = content_box.text_frame
        tf2.word_wrap = True
        for j, bullet in enumerate(bullets):
            p2 = tf2.paragraphs[0] if j == 0 else tf2.add_paragraph()
            p2.text = f"• {bullet}"
            p2.runs[0].font.size = Pt(18)

    # Notes
    if notes:
        slide.notes_slide.notes_text_frame.text = notes

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)

    result = _upload_bytes(buf.read(), filename)
    return {"status": "updated", "slide_count": len(prs.slides), **result}


def get_presentation_summary(filename: str) -> dict:
    """
    Get a high-level summary of a PowerPoint — just titles and bullet counts per slide.
    Faster than read_presentation for large decks.
    """
    _require_pptx()
    item = _find_onedrive_item(filename)
    raw = _download_to_bytes(item["id"])
    prs = Presentation(io.BytesIO(raw))

    summary = []
    for i, slide in enumerate(prs.slides, 1):
        title = ""
        bullet_count = 0
        for shape in slide.shapes:
            if shape.has_text_frame:
                if hasattr(shape, "placeholder_format") and shape.placeholder_format and shape.placeholder_format.idx == 0:
                    title = shape.text_frame.text.strip()
                else:
                    bullet_count += len([l for l in shape.text_frame.text.split("\n") if l.strip()])
        summary.append({
            "slide": i,
            "title": title or f"Slide {i}",
            "bullet_count": bullet_count,
        })

    return {
        "filename": item["name"],
        "slide_count": len(prs.slides),
        "slides": summary,
        "webUrl": item.get("webUrl", ""),
    }
