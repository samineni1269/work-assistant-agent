"""
tools/doc_creator.py — Word Document & PowerPoint Creator
==========================================================
Creates .docx and .pptx files from structured content supplied by the agent.
All created files are stored in ~/work-assistant-docs/ and tracked in a
local SQLite library so they can be browsed in the web UI.

Usage (called by the agent as tools):
    create_word_document(title, sections, filename=None)
    create_presentation(title, slides, filename=None)
    list_documents(doc_type=None)
    delete_document(doc_id)
    get_document_path(doc_id)
"""

import os
import json
import sqlite3
import datetime
from pathlib import Path

# ── Storage
DOCS_DIR = Path.home() / "work-assistant-docs"
DOCS_DIR.mkdir(exist_ok=True)
DB_PATH  = Path.home() / ".work-assistant-docs.db"


# ─────────────────────────────────────────────
# DATABASE
# ─────────────────────────────────────────────

def _get_db() -> sqlite3.Connection:
    conn = sqlite3.connect(str(DB_PATH))
    conn.row_factory = sqlite3.Row
    conn.execute("""
        CREATE TABLE IF NOT EXISTS documents (
            id          INTEGER PRIMARY KEY AUTOINCREMENT,
            title       TEXT    NOT NULL,
            doc_type    TEXT    NOT NULL,   -- 'docx' or 'pptx'
            filename    TEXT    NOT NULL,
            filepath    TEXT    NOT NULL,
            size_bytes  INTEGER DEFAULT 0,
            created_at  TEXT    NOT NULL,
            slide_count INTEGER DEFAULT 0,
            page_count  INTEGER DEFAULT 0
        )
    """)
    conn.commit()
    return conn


def _register(title, doc_type, filename, filepath, slide_count=0, page_count=0):
    conn = _get_db()
    size = Path(filepath).stat().st_size if Path(filepath).exists() else 0
    now  = datetime.datetime.now().isoformat()
    cur  = conn.execute("""
        INSERT INTO documents (title, doc_type, filename, filepath, size_bytes, created_at, slide_count, page_count)
        VALUES (?, ?, ?, ?, ?, ?, ?, ?)
    """, (title, doc_type, filename, str(filepath), size, now, slide_count, page_count))
    conn.commit()
    doc_id = cur.lastrowid
    conn.close()
    return doc_id


def list_documents(doc_type: str = None) -> list[dict]:
    """Return all documents in the library, optionally filtered by type."""
    conn = _get_db()
    if doc_type:
        rows = conn.execute(
            "SELECT * FROM documents WHERE doc_type=? ORDER BY created_at DESC", (doc_type,)
        ).fetchall()
    else:
        rows = conn.execute(
            "SELECT * FROM documents ORDER BY created_at DESC"
        ).fetchall()
    conn.close()
    return [dict(r) for r in rows]


def delete_document(doc_id: int) -> dict:
    """Delete a document from the library and disk."""
    conn = _get_db()
    row = conn.execute("SELECT filepath FROM documents WHERE id=?", (doc_id,)).fetchone()
    if row:
        try:
            Path(row["filepath"]).unlink(missing_ok=True)
        except Exception:
            pass
        conn.execute("DELETE FROM documents WHERE id=?", (doc_id,))
        conn.commit()
    conn.close()
    return {"status": "deleted", "id": doc_id}


def get_document_path(doc_id: int) -> str:
    """Return the filesystem path for a document by ID."""
    conn = _get_db()
    row = conn.execute("SELECT filepath FROM documents WHERE id=?", (doc_id,)).fetchone()
    conn.close()
    return row["filepath"] if row else ""


# ─────────────────────────────────────────────
# WORD DOCUMENT CREATOR
# ─────────────────────────────────────────────

def create_word_document(
    title: str,
    sections: list[dict],
    filename: str = None,
    subtitle: str = "",
) -> dict:
    """
    Create a formatted Word (.docx) document.

    Args:
        title:    Document title (also used as filename if not provided)
        sections: List of section dicts, each with:
                    - heading (str): section heading
                    - content (list[str]): paragraphs or bullet points
                    - style (str): 'paragraph', 'bullets', or 'table'
                    - table_data (list[list[str]]): rows for style='table'
        filename: Optional output filename (without extension)
        subtitle: Optional subtitle shown below the title

    Returns:
        {"status": "created", "id": int, "path": str, "filename": str}
    """
    try:
        from docx import Document
        from docx.shared import Pt, RGBColor, Inches
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        from docx.oxml.ns import qn
        from docx.oxml import OxmlElement
    except ImportError:
        return {"error": "python-docx not installed. Run: pip install python-docx"}

    doc = Document()

    # ── Page margins
    for section in doc.sections:
        section.top_margin    = Inches(1)
        section.bottom_margin = Inches(1)
        section.left_margin   = Inches(1.2)
        section.right_margin  = Inches(1.2)

    # ── Title
    title_para = doc.add_paragraph()
    title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title_para.add_run(title)
    run.bold      = True
    run.font.size = Pt(28)
    run.font.color.rgb = RGBColor(0x1F, 0x4E, 0x79)

    # ── Subtitle
    if subtitle:
        sub_para = doc.add_paragraph()
        sub_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        sub_run = sub_para.add_run(subtitle)
        sub_run.font.size = Pt(14)
        sub_run.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
        sub_run.italic = True

    # ── Date line
    date_para = doc.add_paragraph()
    date_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    date_run = date_para.add_run(datetime.datetime.now().strftime("%B %Y"))
    date_run.font.size = Pt(11)
    date_run.font.color.rgb = RGBColor(0x99, 0x99, 0x99)

    doc.add_paragraph()  # spacer

    page_count = 0

    for sec in sections:
        heading_text = sec.get("heading", "")
        content      = sec.get("content", [])
        style        = sec.get("style", "paragraph")
        table_data   = sec.get("table_data", [])

        if heading_text:
            h = doc.add_heading(heading_text, level=1)
            for run in h.runs:
                run.font.color.rgb = RGBColor(0x2E, 0x75, 0xB6)
            h.paragraph_format.space_before = Pt(12)
            h.paragraph_format.space_after  = Pt(6)

        if style == "bullets":
            for item in content:
                p = doc.add_paragraph(style="List Bullet")
                run = p.add_run(str(item))
                run.font.size = Pt(11)

        elif style == "table" and table_data:
            if len(table_data) > 0:
                cols = max(len(row) for row in table_data)
                table = doc.add_table(rows=len(table_data), cols=cols)
                table.style = "Table Grid"
                for r_idx, row in enumerate(table_data):
                    for c_idx, cell_text in enumerate(row):
                        cell = table.cell(r_idx, c_idx)
                        cell.text = str(cell_text)
                        if r_idx == 0:
                            for run in cell.paragraphs[0].runs:
                                run.bold = True
                                run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                            # Blue header background
                            tc_pr = cell._tc.get_or_add_tcPr()
                            shd = OxmlElement('w:shd')
                            shd.set(qn('w:val'), 'clear')
                            shd.set(qn('w:color'), 'auto')
                            shd.set(qn('w:fill'), '2E75B6')
                            tc_pr.append(shd)
            doc.add_paragraph()

        else:  # paragraph
            for para_text in content:
                p = doc.add_paragraph()
                run = p.add_run(str(para_text))
                run.font.size = Pt(11)
                p.paragraph_format.space_after = Pt(6)

        page_count += max(1, len(content) // 8)

    page_count = max(1, page_count)

    # ── Save
    safe_name = _safe_filename(filename or title)
    out_path  = DOCS_DIR / f"{safe_name}.docx"
    # Avoid overwriting
    counter = 1
    while out_path.exists():
        out_path = DOCS_DIR / f"{safe_name}_{counter}.docx"
        counter += 1

    doc.save(str(out_path))

    doc_id = _register(title, "docx", out_path.name, str(out_path), page_count=page_count)
    return {
        "status":   "created",
        "id":       doc_id,
        "path":     str(out_path),
        "filename": out_path.name,
        "type":     "docx",
        "message":  f"Word document '{out_path.name}' created successfully with {len(sections)} sections.",
    }


# ─────────────────────────────────────────────
# POWERPOINT CREATOR
# ─────────────────────────────────────────────

def create_presentation(
    title: str,
    slides: list[dict],
    filename: str = None,
    theme: str = "blue",
) -> dict:
    """
    Create a formatted PowerPoint (.pptx) presentation.

    Args:
        title:   Presentation title (also used as filename if not provided)
        slides:  List of slide dicts, each with:
                   - title (str): slide title
                   - content (list[str]): bullet points or paragraphs
                   - layout (str): 'title', 'bullets', 'two_col', 'blank'
                   - notes (str): optional speaker notes
                   - col1 (list[str]): left column for two_col layout
                   - col2 (list[str]): right column for two_col layout
        filename: Optional output filename (without extension)
        theme:   'blue' (default), 'dark', or 'minimal'

    Returns:
        {"status": "created", "id": int, "path": str, "filename": str, "slide_count": int}
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        return {"error": "python-pptx not installed. Run: pip install python-pptx"}

    # Theme colours
    themes = {
        "blue":    {"bg": "FFFFFF", "title_bg": "1F4E79", "title_fg": "FFFFFF", "h_fg": "1F4E79", "body_fg": "1A1A2E", "accent": "2E75B6"},
        "dark":    {"bg": "1A1A2E", "title_bg": "0D0D1A", "title_fg": "FFFFFF", "h_fg": "00B0F0", "body_fg": "CCCCCC", "accent": "00B0F0"},
        "minimal": {"bg": "FAFAFA", "title_bg": "222222", "title_fg": "FFFFFF", "h_fg": "222222", "body_fg": "444444", "accent": "888888"},
    }
    T = themes.get(theme, themes["blue"])

    def hex_color(h):
        h = h.lstrip("#")
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # completely blank

    def add_rect(slide, x, y, w, h, color_hex):
        from pptx.util import Inches
        shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))  # MSO_SHAPE_TYPE.RECTANGLE = 1
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_color(color_hex)
        shape.line.fill.background()
        return shape

    def add_text_box(slide, text, x, y, w, h, font_size, bold=False, color_hex="000000", align="left", italic=False, wrap=True):
        from pptx.util import Inches, Pt
        txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = txBox.text_frame
        tf.word_wrap = wrap
        p = tf.paragraphs[0]
        p.text = text
        run = p.runs[0] if p.runs else p.add_run()
        run.text = text
        run.font.size  = Pt(font_size)
        run.font.bold  = bold
        run.font.italic = italic
        run.font.color.rgb = hex_color(color_hex)
        if align == "center":
            p.alignment = PP_ALIGN.CENTER
        elif align == "right":
            p.alignment = PP_ALIGN.RIGHT
        else:
            p.alignment = PP_ALIGN.LEFT
        return txBox

    def add_bullet_box(slide, items, x, y, w, h, font_size=16, color_hex="1A1A2E"):
        from pptx.util import Inches, Pt
        txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"• {item}"
            if p.runs:
                run = p.runs[0]
            else:
                run = p.add_run()
                run.text = f"• {item}"
            run.font.size  = Pt(font_size)
            run.font.color.rgb = hex_color(color_hex)
        return txBox

    # ── SLIDE 1: Title slide
    slide = prs.slides.add_slide(blank_layout)
    add_rect(slide, 0, 0, 13.33, 7.5, T["title_bg"])
    add_rect(slide, 0, 0, 0.2, 7.5, T["accent"])
    add_text_box(slide, title, 0.5, 2.0, 12.5, 1.8, 52, bold=True, color_hex=T["title_fg"])
    add_text_box(slide, f"Created {datetime.datetime.now().strftime('%B %Y')}", 0.5, 4.0, 8, 0.5, 14, color_hex="AAAAAA")

    # ── Content slides
    for sd in slides:
        slide_title   = sd.get("title", "")
        content       = sd.get("content", [])
        layout        = sd.get("layout", "bullets")
        notes_text    = sd.get("notes", "")
        col1          = sd.get("col1", [])
        col2          = sd.get("col2", [])

        slide = prs.slides.add_slide(blank_layout)
        # Background
        add_rect(slide, 0, 0, 13.33, 7.5, T["bg"])
        # Header bar
        add_rect(slide, 0, 0, 13.33, 1.1, T["title_bg"])
        # Title text
        if slide_title:
            add_text_box(slide, slide_title, 0.4, 0.12, 12.5, 0.85, 26, bold=True, color_hex=T["title_fg"])
        # Accent divider
        add_rect(slide, 0.4, 1.18, 12.53, 0.04, T["accent"])

        if layout == "two_col":
            add_bullet_box(slide, col1 or content[:len(content)//2], 0.4, 1.35, 6.0, 5.8, 16, T["body_fg"])
            add_bullet_box(slide, col2 or content[len(content)//2:], 6.7, 1.35, 6.0, 5.8, 16, T["body_fg"])
        elif layout == "blank":
            if content:
                add_text_box(slide, "\n".join(content), 0.4, 1.35, 12.53, 5.8, 16, color_hex=T["body_fg"])
        else:  # bullets (default)
            if content:
                add_bullet_box(slide, content, 0.4, 1.35, 12.53, 5.8, 16, T["body_fg"])

        # Speaker notes
        if notes_text:
            slide.notes_slide.notes_text_frame.text = notes_text

    # ── Closing slide
    closing = prs.slides.add_slide(blank_layout)
    add_rect(closing, 0, 0, 13.33, 7.5, T["title_bg"])
    add_rect(closing, 0, 0, 0.2, 7.5, T["accent"])
    add_text_box(closing, "Thank You", 0.5, 2.8, 12.5, 1.4, 48, bold=True, color_hex=T["title_fg"], align="center")

    # ── Save
    safe_name  = _safe_filename(filename or title)
    out_path   = DOCS_DIR / f"{safe_name}.pptx"
    counter = 1
    while out_path.exists():
        out_path = DOCS_DIR / f"{safe_name}_{counter}.pptx"
        counter += 1

    prs.save(str(out_path))
    slide_count = len(prs.slides)

    doc_id = _register(title, "pptx", out_path.name, str(out_path), slide_count=slide_count)
    return {
        "status":      "created",
        "id":          doc_id,
        "path":        str(out_path),
        "filename":    out_path.name,
        "type":        "pptx",
        "slide_count": slide_count,
        "message":     f"Presentation '{out_path.name}' created with {slide_count} slides.",
    }


# ─────────────────────────────────────────────
# HELPERS
# ─────────────────────────────────────────────

def _safe_filename(name: str) -> str:
    """Convert a title to a safe filename."""
    import re
    safe = re.sub(r'[^\w\s\-]', '', name)
    safe = re.sub(r'\s+', '_', safe.strip())
    return safe[:60] or "document"
